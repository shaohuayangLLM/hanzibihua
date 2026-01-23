#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT文件解析器
用于提取一年级数学期末测评PPT中的题目信息
"""

from pptx import Presentation
import json
import os

def extract_text_from_shape(shape):
    """从形状中提取文本"""
    text = ""
    if hasattr(shape, "text"):
        text = shape.text
    return text.strip()

def extract_table_data(table):
    """从表格中提取数据"""
    data = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        data.append(row_data)
    return data

def parse_ppt_file(file_path):
    """解析PPT文件，提取所有内容"""
    try:
        prs = Presentation(file_path)
        slides_data = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_info = {
                "slide_number": slide_idx + 1,
                "title": "",
                "content": [],
                "tables": [],
                "images": []
            }

            # 提取幻灯片中的所有形状
            for shape in slide.shapes:
                # 提取文本
                text = extract_text_from_shape(shape)
                if text:
                    # 判断是否是标题（通常是第一个文本框或较大字体的文本）
                    if not slide_info["title"] and len(text) < 50:
                        slide_info["title"] = text
                    else:
                        slide_info["content"].append(text)

                # 提取表格
                if shape.shape_type == 19:  # SHAPE_TYPE.TABLE
                    table_data = extract_table_data(shape.table)
                    slide_info["tables"].append(table_data)

                # 提取图片信息
                if shape.shape_type == 13:  # SHAPE_TYPE.PICTURE
                    slide_info["images"].append(f"图片_{len(slide_info['images']) + 1}")

            slides_data.append(slide_info)

        return slides_data
    except Exception as e:
        print(f"解析PPT文件时出错: {e}")
        return None

def analyze_math_problems(slides_data):
    """分析数学题目"""
    analysis = {
        "question_types": [],
        "examples": [],
        "knowledge_points": [],
        "difficulty_levels": []
    }

    current_section = None

    for slide in slides_data:
        content = " ".join(slide["content"])
        title = slide["title"]

        # 识别题型
        if any(keyword in title for keyword in ["填空", "填一填"]):
            current_section = "填空题"
            analysis["question_types"].append({
                "type": "填空题",
                "description": "根据题目要求填写适当的数字或答案",
                "slide": slide["slide_number"]
            })
        elif any(keyword in title for keyword in ["比一比", "比较"]):
            current_section = "比较题"
            analysis["question_types"].append({
                "type": "比较题",
                "description": "比较大小、多少、长短等",
                "slide": slide["slide_number"]
            })
        elif any(keyword in title for keyword in ["数一数", "计数"]):
            current_section = "计数题"
            analysis["question_types"].append({
                "type": "计数题",
                "description": "数物品的数量",
                "slide": slide["slide_number"]
            })
        elif any(keyword in title for keyword in ["算一算", "计算"]):
            current_section = "计算题"
            analysis["question_types"].append({
                "type": "计算题",
                "description": "进行加减法运算",
                "slide": slide["slide_number"]
            })
        elif any(keyword in title for keyword in ["连一连", "连线"]):
            current_section = "连线题"
            analysis["question_types"].append({
                "type": "连线题",
                "description": "将相关内容用线连接",
                "slide": slide["slide_number"]
            })
        elif any(keyword in title for keyword in ["看图", "观察"]):
            current_section = "图形题"
            analysis["question_types"].append({
                "type": "图形题",
                "description": "观察图形并回答问题",
                "slide": slide["slide_number"]
            })

        # 提取题目示例
        if slide["content"]:
            analysis["examples"].append({
                "slide": slide["slide_number"],
                "section": current_section,
                "title": title,
                "content": slide["content"][:5],  # 限制内容长度
                "has_images": len(slide["images"]) > 0,
                "has_tables": len(slide["tables"]) > 0
            })

        # 提取表格数据（可能包含数字序列、算式等）
        if slide["tables"]:
            for table in slide["tables"]:
                analysis["examples"].append({
                    "slide": slide["slide_number"],
                    "section": current_section or "表格题",
                    "type": "table",
                    "data": table
                })

    return analysis

def main():
    """主函数"""
    base_path = "/Users/ysh/Manual Library/ClaudeCode/portfolio-2025/K12-Education/一年级数学上"
    ppt_files = [
        "期末综合素质测评(一).pptx",
        "期末综合素质测评(二).pptx"
    ]

    all_results = {}

    for ppt_file in ppt_files:
        file_path = os.path.join(base_path, ppt_file)
        print(f"\n正在解析: {ppt_file}")

        slides_data = parse_ppt_file(file_path)
        if slides_data:
            analysis = analyze_math_problems(slides_data)
            all_results[ppt_file] = {
                "slides_count": len(slides_data),
                "slides_data": slides_data,
                "analysis": analysis
            }
            print(f"  ✓ 成功解析 {len(slides_data)} 张幻灯片")
        else:
            print(f"  ✗ 解析失败")

    # 保存结果为JSON
    output_path = "/Users/ysh/Manual Library/ClaudeCode/portfolio-2025/K12-Education/src/data/math/generators/ppt_analysis_result.json"
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(all_results, f, ensure_ascii=False, indent=2)

    print(f"\n结果已保存到: {output_path}")

    # 打印摘要
    print("\n=== 解析摘要 ===")
    for ppt_file, data in all_results.items():
        print(f"\n{ppt_file}:")
        print(f"  幻灯片数量: {data['slides_count']}")
        if data['analysis']['question_types']:
            print(f"  识别的题型:")
            for qt in data['analysis']['question_types']:
                print(f"    - {qt['type']}: {qt.get('description', '')}")

if __name__ == "__main__":
    main()
